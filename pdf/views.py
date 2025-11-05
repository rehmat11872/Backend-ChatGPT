import os
from rest_framework import generics
from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status
from PyPDF2 import PdfReader, PdfWriter
from rest_framework.permissions import IsAuthenticated, AllowAny
from rest_framework.parsers import MultiPartParser, FormParser
from django.http import FileResponse
from django.contrib.auth.hashers import make_password
from .models import OcrPdf, ProtectedPDF, WordToPdfConversion, WordToPdf, OrganizedPdf, MergedPDF, CompressedPDF, SplitPDF, UnlockPdf, PDFFormatConversion, StampPdf
from django.core.files.base import ContentFile
from .utils import convert_pdf_to_image, create_zip_file, stamp_pdf_with_text, pdf_to_ocr
from .serializers import OcrPdfSerializer, ProtectedPDFSerializer, MergedPDFSerializer, CompressedPDFSerializer, SplitPDFSerializer, StampPdfSerializer, WordToPdfConversionSerializer, OrganizedPdfSerializer, UnlockPdfSerializer, ProtectPDFRequestSerializer, MergePDFRequestSerializer, CompressPDFRequestSerializer, SplitPDFRequestSerializer, WordToPdfRequestSerializer, OrganizePDFRequestSerializer, UnlockPDFRequestSerializer, StampPDFRequestSerializer, OcrPDFRequestSerializer, PDFToFormatRequestSerializer, PDFFormatConversionSerializer
from drf_spectacular.utils import extend_schema, inline_serializer, OpenApiExample
from rest_framework import serializers as drf_serializers


@extend_schema(tags=['PDF Operations'])
class ProtectPDFView(APIView):
    permission_classes = [AllowAny]
    parser_classes = (MultiPartParser, FormParser)
    serializer_class = ProtectPDFRequestSerializer

    @extend_schema(
        request=ProtectPDFRequestSerializer,
        responses=inline_serializer(
            name='ProtectPDFResponse',
            fields={
                'message': drf_serializers.CharField(),
                'split_pdf': inline_serializer(
                    name='ProtectPDFData',
                    fields={
                        'id': drf_serializers.IntegerField(),
                        'user': drf_serializers.IntegerField(),
                        'created_at': drf_serializers.DateTimeField(),
                        'protected_file': drf_serializers.URLField(),
                    }
                ),
            }
        ),
        examples=[
            OpenApiExample(
                'Protect PDF Example',
                value={
                    'message': 'PDF protection completed',
                    'split_pdf': {
                        'id': 1,
                        'user': 5,
                        'created_at': '2024-01-01T00:00:00Z',
                        'protected_file': 'http://localhost:8000/media/protected/file.pdf'
                    }
                }
            )
        ]
    )
    def post(self, request, format=None):
        input_file = request.data.get('input_pdf', None)
        pdf_password = request.data.get('pdf_password', None)
        
        # Get permissions
        allow_printing = request.data.get('allow_printing', 'true').lower() == 'true'
        allow_copying = request.data.get('allow_copying', 'true').lower() == 'true'
        allow_editing = request.data.get('allow_editing', 'false').lower() == 'true'
        allow_comments = request.data.get('allow_comments', 'true').lower() == 'true'
        allow_form_filling = request.data.get('allow_form_filling', 'true').lower() == 'true'
        allow_document_assembly = request.data.get('allow_document_assembly', 'false').lower() == 'true'

        if not input_file or not pdf_password:
            return Response({'error': 'Incomplete data provided.'}, status=status.HTTP_400_BAD_REQUEST)

        try:
            # Create user if not authenticated
            if request.user.is_authenticated:
                user = request.user
            else:
                from accounts.models import User
                user, created = User.objects.get_or_create(
                    email='test@example.com',
                    defaults={'password': make_password('testpass123')}
                )
            
            # Protect PDF with permissions
            reader = PdfReader(input_file)
            writer = PdfWriter()
            
            for page in reader.pages:
                writer.add_page(page)
            
            # Set permissions based on user input
            writer.encrypt(
                user_password=pdf_password,
                owner_password=pdf_password + '_owner',
                use_128bit=True,
                permissions_flag=(
                    (4 if allow_printing else 0) |
                    (16 if allow_copying else 0) |
                    (32 if allow_editing else 0) |
                    (64 if allow_comments else 0) |
                    (256 if allow_form_filling else 0) |
                    (1024 if allow_document_assembly else 0)
                )
            )
            
            # Save protected PDF
            from io import BytesIO
            buffer = BytesIO()
            writer.write(buffer)
            buffer.seek(0)
            
            protected_pdf = ProtectedPDF(user=user)
            protected_pdf.protected_file.save(
                f'protected_{input_file.name}',
                ContentFile(buffer.getvalue())
            )
            protected_pdf.save()
            
            # Get full URL
            protocol = 'https' if request.is_secure() else 'http'
            file_url = f'{protocol}://{request.get_host()}{protected_pdf.protected_file.url}'
            
            response_data = {
                'message': 'PDF protection completed',
                'split_pdf': {
                    'id': protected_pdf.id,
                    'user': user.id,
                    'created_at': protected_pdf.created_at.isoformat(),
                    'protected_file': file_url
                }
            }
            return Response(response_data)
        except Exception as e:
            return Response({'error': f'An error occurred: {str(e)}'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@extend_schema(tags=['PDF Operations'])
class DownloadProtectedPDFView(APIView):
    permission_classes = [AllowAny]

    def get(self, request, pdf_id, format=None):
        try:
            protected_pdf = ProtectedPDF.objects.get(id=pdf_id)
            return FileResponse(
                protected_pdf.protected_file.open('rb'),
                as_attachment=True,
                filename=f'protected_{protected_pdf.id}.pdf'
            )
        except ProtectedPDF.DoesNotExist:
            return Response({'error': 'Protected PDF not found'}, status=404)



@extend_schema(exclude=True)
class ProtectedPDFDeleteView(generics.DestroyAPIView):
    queryset = ProtectedPDF.objects.all()
    serializer_class = ProtectedPDFSerializer
    permission_classes = [IsAuthenticated]


@extend_schema(tags=['PDF Operations'])
class MergePDFView(APIView):
    permission_classes = [IsAuthenticated]
    parser_classes = (MultiPartParser, FormParser)
    serializer_class = MergePDFRequestSerializer

    @extend_schema(
        request=MergePDFRequestSerializer,
        responses=inline_serializer(
            name='MergePDFResponse',
            fields={
                'message': drf_serializers.CharField(),
                'merged_data': inline_serializer(
                    name='MergePDFData',
                    fields={
                        'id': drf_serializers.IntegerField(),
                        'user': drf_serializers.IntegerField(),
                        'created_at': drf_serializers.DateTimeField(),
                        'merged_file': drf_serializers.URLField(),
                    }
                ),
            }
        ),
        examples=[
            OpenApiExample(
                'Merge PDF Example',
                value={
                    'message': 'PDFs merged and saved successfully',
                    'merged_data': {
                        'id': 1,
                        'user': 5,
                        'created_at': '2024-01-01T00:00:00Z',
                        'merged_file': 'http://localhost:8000/media/merged/output.pdf'
                    }
                }
            )
        ]
    )
    def post(self, request, format=None):
        pdf_files = request.FILES.getlist('pdf_files')

        if len(pdf_files) < 2:
            return Response({'error': 'At least two PDFs are required for merging.'}, status=status.HTTP_400_BAD_REQUEST)

        try:
            user = request.user
            
            # Merge PDFs
            writer = PdfWriter()
            for pdf_file in pdf_files:
                try:
                    reader = PdfReader(pdf_file)
                    # Handle encrypted PDFs
                    if reader.is_encrypted:
                        return Response({'error': f'PDF file "{pdf_file.name}" is encrypted. Please unlock it first.'}, status=status.HTTP_400_BAD_REQUEST)
                    
                    for page in reader.pages:
                        writer.add_page(page)
                except Exception as pdf_error:
                    return Response({'error': f'Error processing PDF "{pdf_file.name}": {str(pdf_error)}'}, status=status.HTTP_400_BAD_REQUEST)
            
            # Save merged PDF
            from io import BytesIO
            buffer = BytesIO()
            writer.write(buffer)
            buffer.seek(0)
            
            merged_pdf = MergedPDF(user=user)
            merged_pdf.merged_file.save(
                f'merged_{len(pdf_files)}_files.pdf',
                ContentFile(buffer.getvalue())
            )
            merged_pdf.save()
            
            # Get full URL
            protocol = 'https' if request.is_secure() else 'http'
            file_url = f'{protocol}://{request.get_host()}{merged_pdf.merged_file.url}'
            
            response_data = {
                'message': 'PDFs merged and saved successfully',
                'merged_data': {
                    'id': merged_pdf.id,
                    'user': user.id,
                    'created_at': merged_pdf.created_at.isoformat(),
                    'merged_file': file_url
                },
            }
            return Response(response_data)

        except Exception as e:
            return Response({'error': f'An error occurred: {str(e)}'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@extend_schema(tags=['PDF Operations'])
class DownloadMergedPDFView(APIView):
    permission_classes = [AllowAny]

    def get(self, request, pdf_id, format=None):
        try:
            merged_pdf = MergedPDF.objects.get(id=pdf_id)
            return FileResponse(
                merged_pdf.merged_file.open('rb'),
                as_attachment=True,
                filename=f'merged_{merged_pdf.id}.pdf'
            )
        except MergedPDF.DoesNotExist:
            return Response({'error': 'Merged PDF not found'}, status=404)


@extend_schema(exclude=True)
class MergePDFDeleteView(generics.DestroyAPIView):
    queryset = MergedPDF.objects.all()
    serializer_class = MergedPDFSerializer
    permission_classes = [IsAuthenticated]

@extend_schema(tags=['PDF Operations'])
class CompressPDFView(APIView):
    permission_classes = [AllowAny]
    parser_classes = (MultiPartParser, FormParser)
    serializer_class = CompressPDFRequestSerializer

    @extend_schema(
        request=CompressPDFRequestSerializer,
        responses=inline_serializer(
            name='CompressPDFResponse',
            fields={
                'message': drf_serializers.CharField(),
                'split_pdf': inline_serializer(
                    name='CompressPDFData',
                    fields={
                        'id': drf_serializers.IntegerField(),
                        'user': drf_serializers.IntegerField(),
                        'created_at': drf_serializers.DateTimeField(),
                        'compressed_file': drf_serializers.URLField(),
                    }
                ),
            }
        ),
        examples=[
            OpenApiExample(
                'Compress PDF Example',
                value={
                    'message': 'PDF compression completed',
                    'split_pdf': {
                        'id': 1,
                        'user': 5,
                        'created_at': '2024-01-01T00:00:00Z',
                        'compressed_file': 'http://localhost:8000/media/compressed/output.pdf'
                    }
                }
            )
        ]
    )
    def post(self, request, format=None):
        input_pdf = request.FILES.get('input_pdf', None)
        compression_quality = request.data.get('compression_quality', 'recommended')

        if not input_pdf:
            return Response({'error': 'No input PDF file provided.'}, status=status.HTTP_400_BAD_REQUEST)

        try:
            # Create user if not authenticated
            if request.user.is_authenticated:
                user = request.user
            else:
                from accounts.models import User
                user, created = User.objects.get_or_create(
                    email='test@example.com',
                    defaults={'password': make_password('testpass123')}
                )
            
            # Compress PDF (simple copy for now)
            reader = PdfReader(input_pdf)
            writer = PdfWriter()
            
            for page in reader.pages:
                writer.add_page(page)
            
            # Save compressed PDF
            from io import BytesIO
            buffer = BytesIO()
            writer.write(buffer)
            buffer.seek(0)
            
            compressed_pdf = CompressedPDF(user=user)
            compressed_pdf.compressed_file.save(
                f'compressed_{compression_quality}.pdf',
                ContentFile(buffer.getvalue())
            )
            compressed_pdf.save()
            
            # Get full URL
            protocol = 'https' if request.is_secure() else 'http'
            file_url = f'{protocol}://{request.get_host()}{compressed_pdf.compressed_file.url}'
            
            response_data = {
                'message': 'PDF compression completed',
                'split_pdf': {
                    'id': compressed_pdf.id,
                    'user': user.id,
                    'created_at': compressed_pdf.created_at.isoformat(),
                    'compressed_file': file_url
                },
            }
            return Response(response_data)

        except Exception as e:
            return Response({'error': f'An error occurred: {str(e)}'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@extend_schema(tags=['PDF Operations'])
class DownloadCompressedPDFView(APIView):
    permission_classes = [AllowAny]

    def get(self, request, pdf_id, format=None):
        try:
            compressed_pdf = CompressedPDF.objects.get(id=pdf_id)
            return FileResponse(
                compressed_pdf.compressed_file.open('rb'),
                as_attachment=True,
                filename=f'compressed_{compressed_pdf.id}.pdf'
            )
        except CompressedPDF.DoesNotExist:
            return Response({'error': 'Compressed PDF not found'}, status=404)


@extend_schema(exclude=True)
class CompressPDFDeleteView(generics.DestroyAPIView):
    queryset = CompressedPDF.objects.all()
    serializer_class = CompressedPDFSerializer
    permission_classes = [IsAuthenticated]



from io import BytesIO
import json


@extend_schema(tags=['PDF Operations'])
class SplitPDFView(APIView):
    permission_classes = [IsAuthenticated]
    parser_classes = (MultiPartParser, FormParser)
    serializer_class = SplitPDFRequestSerializer

    @extend_schema(
        request=SplitPDFRequestSerializer,
        responses=inline_serializer(
            name='SplitPDFResponse',
            fields={
                'message': drf_serializers.CharField(),
                'split_pdf': inline_serializer(
                    name='SplitPDFData',
                    fields={
                        'id': drf_serializers.IntegerField(),
                        'user': drf_serializers.IntegerField(),
                        'created_at': drf_serializers.DateTimeField(),
                        'split_pdf': drf_serializers.URLField(),
                    }
                ),
            }
        ),
        examples=[
            OpenApiExample(
                'Split PDF Example',
                value={
                    'message': 'PDF splitting completed',
                    'split_pdf': {
                        'id': 1,
                        'user': 5,
                        'created_at': '2024-01-01T00:00:00Z',
                        'split_pdf': 'http://localhost:8000/media/split/output_part1.pdf'
                    }
                }
            )
        ]
    )
    def post(self, request, format=None):
        input_pdf = request.FILES.get('input_pdf')
        if not input_pdf:
            return Response({'error': 'No input PDF file provided.'}, status=status.HTTP_400_BAD_REQUEST)

        user = request.user
        reader = PdfReader(input_pdf)

        # --- Extract split options from frontend ---
        split_mode = request.data.get('split_mode', 'range')
        range_mode = request.data.get('range_mode', 'custom')
        ranges_json = request.data.get('ranges', '[]')
        extract_mode = request.data.get('extract_mode', 'all')
        pages_to_extract = request.data.get('pages_to_extract', '')
        max_file_size = request.data.get('max_file_size', '')
        merge_ranges = request.data.get('merge_ranges', 'false') == 'true'
        merge_extracted = request.data.get('merge_extracted', 'false') == 'true'

        try:
            if split_mode == "range":
                if isinstance(ranges_json, str):
                    ranges_json = ranges_json.strip().rstrip("'\"")
                    ranges = json.loads(ranges_json)
                else:
                    ranges = ranges_json
            else:
                ranges = []
        except Exception as e:
            return Response({'error': f'Invalid ranges JSON format: {str(e)}'}, status=status.HTTP_400_BAD_REQUEST)

        saved_files = []

        try:
            # --- Split by mode ---
            if split_mode == "range":
                saved_files = self._split_by_range(reader, user, ranges, range_mode, merge_ranges, request)
            elif split_mode == "pages":
                saved_files = self._split_by_pages(reader, user, extract_mode, pages_to_extract, merge_extracted, request)
            elif split_mode == "size":
                saved_files = self._split_by_size(reader, user, max_file_size, request)
            else:
                return Response({'error': 'Invalid split mode'}, status=status.HTTP_400_BAD_REQUEST)

            # Return format expected by frontend with dynamic message
            if not saved_files:
                return Response({'error': 'No files created'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
            
            files_count = len(saved_files)
            
            # If multiple files, create ZIP
            if files_count > 1:
                zip_file = self._create_zip_from_files(saved_files, user, request)
                message = f"PDF split completed. {files_count} PDFs created and zipped."
                
                response_data = {
                    'message': message,
                    'split_pdf': {
                        'id': zip_file['id'],
                        'user': zip_file['user'],
                        'created_at': zip_file['created_at'],
                        'split_pdf': zip_file['split_pdf']
                    },
                    'total_files_created': files_count,
                    'is_zip': True
                }
            else:
                # Single file, return as normal
                first_file = saved_files[0]
                message = f"PDF split completed. 1 PDF created."
                
                response_data = {
                    'message': message,
                    'split_pdf': {
                        'id': first_file['id'],
                        'user': first_file['user'],
                        'created_at': first_file['created_at'],
                        'split_pdf': first_file['split_pdf']
                    },
                    'total_files_created': files_count,
                    'is_zip': False
                }
            
            return Response(response_data, status=status.HTTP_200_OK)

        except Exception as e:
            return Response({'error': f'An error occurred: {str(e)}'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

    # ----------------------------
    # Split Helper Methods
    # ----------------------------
    def _split_by_range(self, reader, user, ranges, range_mode, merge_ranges, request):
        files = []
        protocol = 'https' if request.is_secure() else 'http'
        
        print(f"DEBUG _split_by_range: ranges={ranges}, range_mode={range_mode}, merge_ranges={merge_ranges}")

        if not ranges:
            print(f"DEBUG: No ranges provided, returning empty list")
            return []

        if merge_ranges:
            # Merge all ranges into a single file
            writer = PdfWriter()
            for r in ranges:
                start, end = int(r.get('from', 1)), int(r.get('to', len(reader.pages)))
                for i in range(start - 1, min(end, len(reader.pages))):
                    writer.add_page(reader.pages[i])
            buffer = BytesIO()
            writer.write(buffer)
            buffer.seek(0)
            split_pdf = SplitPDF(user=user, start_page=1, end_page=len(reader.pages))
            split_pdf.split_pdf.save(f'split_merged.pdf', ContentFile(buffer.getvalue()))
            split_pdf.save()
            files.append({
                'id': split_pdf.id,
                'user': user.id,
                'created_at': split_pdf.created_at.isoformat(),
                'split_pdf': f'{protocol}://{request.get_host()}{split_pdf.split_pdf.url}'
            })
            return files

        # Create individual files per range
        for idx, r in enumerate(ranges, start=1):
            start, end = int(r.get('from', 1)), int(r.get('to', len(reader.pages)))
            writer = PdfWriter()
            for i in range(start - 1, min(end, len(reader.pages))):
                writer.add_page(reader.pages[i])
            buffer = BytesIO()
            writer.write(buffer)
            buffer.seek(0)
            split_pdf = SplitPDF(user=user, start_page=start, end_page=end)
            split_pdf.split_pdf.save(f'split_part_{idx}.pdf', ContentFile(buffer.getvalue()))
            split_pdf.save()
            files.append({
                'id': split_pdf.id,
                'user': user.id,
                'created_at': split_pdf.created_at.isoformat(),
                'split_pdf': f'{protocol}://{request.get_host()}{split_pdf.split_pdf.url}'
            })
        return files

    def _split_by_pages(self, reader, user, extract_mode, pages_to_extract, merge_extracted, request):
        files = []
        protocol = 'https' if request.is_secure() else 'http'
        total_pages = len(reader.pages)

        if extract_mode == "all":
            for i in range(total_pages):
                writer = PdfWriter()
                writer.add_page(reader.pages[i])
                buffer = BytesIO()
                writer.write(buffer)
                buffer.seek(0)
                split_pdf = SplitPDF(user=user, start_page=i+1, end_page=i+1)
                split_pdf.split_pdf.save(f'page_{i + 1}.pdf', ContentFile(buffer.getvalue()))
                split_pdf.save()
                files.append({
                    'id': split_pdf.id,
                    'user': user.id,
                    'created_at': split_pdf.created_at.isoformat(),
                    'split_pdf': f'{protocol}://{request.get_host()}{split_pdf.split_pdf.url}'
                })
            return files

        # If specific pages provided (e.g., "1,6,8")
        pages = self._parse_pages(pages_to_extract, total_pages)
        if not pages:
            return files
            
        if merge_extracted:
            # Merge selected pages into one PDF
            writer = PdfWriter()
            for p in pages:
                writer.add_page(reader.pages[p - 1])
            buffer = BytesIO()
            writer.write(buffer)
            buffer.seek(0)
            split_pdf = SplitPDF(user=user, start_page=min(pages), end_page=max(pages))
            split_pdf.split_pdf.save('extracted_pages.pdf', ContentFile(buffer.getvalue()))
            split_pdf.save()
            return [{
                'id': split_pdf.id,
                'user': user.id,
                'created_at': split_pdf.created_at.isoformat(),
                'split_pdf': f'{protocol}://{request.get_host()}{split_pdf.split_pdf.url}'
            }]
        else:
            # Create separate PDF for each selected page
            for p in pages:
                writer = PdfWriter()
                writer.add_page(reader.pages[p - 1])
                buffer = BytesIO()
                writer.write(buffer)
                buffer.seek(0)
                split_pdf = SplitPDF(user=user, start_page=p, end_page=p)
                split_pdf.split_pdf.save(f'page_{p}.pdf', ContentFile(buffer.getvalue()))
                split_pdf.save()
                files.append({
                    'id': split_pdf.id,
                    'user': user.id,
                    'created_at': split_pdf.created_at.isoformat(),
                    'split_pdf': f'{protocol}://{request.get_host()}{split_pdf.split_pdf.url}'
                })
            return files

    def _split_by_size(self, reader, user, max_file_size, request):
        files = []
        protocol = 'https' if request.is_secure() else 'http'
        try:
            max_size_kb = float(max_file_size or 500)
        except ValueError:
            max_size_kb = 500

        writer = PdfWriter()
        part = 1
        buffer = BytesIO()

        for page in reader.pages:
            writer.add_page(page)
            temp_buffer = BytesIO()
            writer.write(temp_buffer)
            if temp_buffer.tell() / 1024 > max_size_kb:
                # save current writer and start new
                writer.remove_page(-1)
                buffer = BytesIO()
                writer.write(buffer)
                buffer.seek(0)
                split_pdf = SplitPDF(user=user, start_page=1, end_page=1)  # Will be updated with actual pages
                split_pdf.split_pdf.save(f'size_part_{part}.pdf', ContentFile(buffer.getvalue()))
                split_pdf.save()
                files.append({
                    'id': split_pdf.id,
                    'user': user.id,
                    'created_at': split_pdf.created_at.isoformat(),
                    'split_pdf': f'{protocol}://{request.get_host()}{split_pdf.split_pdf.url}'
                })
                part += 1
                writer = PdfWriter()
                writer.add_page(page)
        # Save the last one
        buffer = BytesIO()
        writer.write(buffer)
        buffer.seek(0)
        split_pdf = SplitPDF(user=user, start_page=1, end_page=1)  # Will be updated with actual pages
        split_pdf.split_pdf.save(f'size_part_{part}.pdf', ContentFile(buffer.getvalue()))
        split_pdf.save()
        files.append({
            'id': split_pdf.id,
            'user': user.id,
            'created_at': split_pdf.created_at.isoformat(),
            'split_pdf': f'{protocol}://{request.get_host()}{split_pdf.split_pdf.url}'
        })
        return files

    def _parse_pages(self, pages_str, total_pages):
        pages = set()
        if not pages_str:
            return []
        
        for part in pages_str.split(','):
            part = part.strip()
            if '-' in part:
                # Handle ranges like "3-5"
                start, end = map(int, part.split('-'))
                pages.update(range(start, min(end + 1, total_pages + 1)))
            else:
                # Handle individual pages like "1", "6", "8"
                p = int(part)
                if 1 <= p <= total_pages:
                    pages.add(p)
        return sorted(pages)
    
    def _create_zip_from_files(self, saved_files, user, request):
        import zipfile
        from io import BytesIO
        
        zip_buffer = BytesIO()
        
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            for i, file_data in enumerate(saved_files, 1):
                # Get the actual file from database
                try:
                    split_pdf = SplitPDF.objects.get(id=file_data['id'])
                    if split_pdf.split_pdf and split_pdf.split_pdf.name:
                        # Read file content
                        split_pdf.split_pdf.seek(0)
                        file_content = split_pdf.split_pdf.read()
                        
                        # Add to ZIP with descriptive name
                        filename = f"page_{i}.pdf"
                        zip_file.writestr(filename, file_content)
                except Exception as e:
                    continue
        
        zip_buffer.seek(0)
        
        # Save ZIP as a new SplitPDF entry
        zip_split_pdf = SplitPDF(user=user, start_page=1, end_page=len(saved_files))
        zip_split_pdf.split_pdf.save(
            'split_pages.zip',
            ContentFile(zip_buffer.getvalue())
        )
        zip_split_pdf.save()
        
        protocol = 'https' if request.is_secure() else 'http'
        
        return {
            'id': zip_split_pdf.id,
            'user': user.id,
            'created_at': zip_split_pdf.created_at.isoformat(),
            'split_pdf': f'{protocol}://{request.get_host()}{zip_split_pdf.split_pdf.url}'
        }


@extend_schema(exclude=True)
class SplitPDFDeleteView(generics.DestroyAPIView):
    queryset = SplitPDF.objects.all()
    serializer_class = SplitPDFSerializer
    permission_classes = [IsAuthenticated]








# name should change to OtherToPdf
@extend_schema(tags=['PDF Operations'])
class WordToPdfConversionView(APIView):
    permission_classes = [IsAuthenticated]
    parser_classes = (MultiPartParser, FormParser)
    serializer_class = WordToPdfRequestSerializer

    @extend_schema(
        request=WordToPdfRequestSerializer,
        responses=inline_serializer(
            name='WordToPdfResponse',
            fields={
                'message': drf_serializers.CharField(),
                'data': inline_serializer(
                    name='WordToPdfData',
                    fields={
                        'id': drf_serializers.IntegerField(),
                        'user': drf_serializers.IntegerField(),
                        'created_at': drf_serializers.DateTimeField(),
                        'pdf': drf_serializers.FileField(),
                    }
                ),
            }
        ),
        examples=[
            OpenApiExample(
                'Word to PDF Example',
                value={
                    'message': 'PDF pages stamped successfully.',
                    'data': {
                        'id': 1,
                        'user': 5,
                        'created_at': '2024-01-01T00:00:00Z',
                        'pdf': 'http://localhost:8000/media/converted/output.pdf'
                    }
                }
            )
        ]
    )
    def post(self, request, format=None):
        input_files = request.FILES.getlist('input_files')

        if not input_files:
            return Response({'error': 'No input files provided.'}, status=status.HTTP_400_BAD_REQUEST)

        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            from docx import Document
            import tempfile
            from django.core.files.base import ContentFile
            from io import BytesIO
            
            conversion_instance = WordToPdfConversion(user=request.user)
            conversion_instance.save()
            
            for input_file in input_files:
                if input_file.name.endswith('.docx'):
                    # Extract text from DOCX
                    doc = Document(input_file)
                    text_content = '\n'.join([p.text for p in doc.paragraphs])
                    
                    # Create PDF with ReportLab
                    buffer = BytesIO()
                    p = canvas.Canvas(buffer, pagesize=letter)
                    width, height = letter
                    
                    # Simple text rendering
                    y = height - 50
                    for line in text_content.split('\n'):
                        if y < 50:
                            p.showPage()
                            y = height - 50
                        p.drawString(50, y, line[:80])  # Limit line length
                        y -= 20
                    
                    p.save()
                    pdf_content = buffer.getvalue()
                    buffer.close()
                    
                    # Save to model
                    word_to_pdf_instance = WordToPdf()
                    filename = f"converted_{input_file.name.split('.')[0]}.pdf"
                    word_to_pdf_instance.word_to_pdf.save(filename, ContentFile(pdf_content))
                    word_to_pdf_instance.save()
                    
                    conversion_instance.word_to_pdfs.add(word_to_pdf_instance)
            
            conversion_instance.save()
            
            serializer = WordToPdfConversionSerializer(conversion_instance, context={'request': request})
            return Response({'message': 'Word to PDF conversion completed.', 'conversion_data': serializer.data})
            
        except Exception as e:
            return Response({'error': f'An error occurred: {str(e)}'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)







@extend_schema(tags=['PDF Operations'])
class DownloadSplitPDFView(APIView):
    permission_classes = [AllowAny]

    def get(self, request, pdf_id, format=None):
        try:
            split_pdf = SplitPDF.objects.get(id=pdf_id)
            return FileResponse(
                split_pdf.split_pdf.open('rb'),
                as_attachment=True,
                filename=f'split_{split_pdf.id}.pdf'
            )
        except SplitPDF.DoesNotExist:
            return Response({'error': 'Split PDF not found'}, status=404)


@extend_schema(tags=['PDF Operations'])
class OrganizePDFView(APIView):
    permission_classes = [IsAuthenticated]
    parser_classes = (MultiPartParser, FormParser)
    serializer_class = OrganizePDFRequestSerializer

    @extend_schema(
        request=OrganizePDFRequestSerializer,
        responses=inline_serializer(
            name='OrganizePDFResponse',
            fields={
                'message': drf_serializers.CharField(),
                'organized_data': inline_serializer(
                    name='OrganizePDFData',
                    fields={
                        'id': drf_serializers.IntegerField(),
                        'user': drf_serializers.IntegerField(),
                        'created_at': drf_serializers.DateTimeField(),
                        'organize_pdf': drf_serializers.URLField(),
                    }
                ),
            }
        ),
        examples=[
            OpenApiExample(
                'Organize PDF Example',
                value={
                    'message': 'PDF pages organized successfully.',
                    'organized_data': {
                        'id': 1,
                        'user': 5,
                        'created_at': '2024-01-01T00:00:00Z',
                        'organize_pdf': 'http://localhost:8000/media/organized/output.pdf'
                    }
                }
            )
        ]
    )
    def post(self, request, format=None):
        input_pdf = request.FILES.get('input_pdf', None)
        user_order = request.data.get('user_order', '')
        delete_pages = request.data.get('delete_pages', '')
        
        if not input_pdf:
            return Response({'error': 'No input PDF file provided.'}, status=status.HTTP_400_BAD_REQUEST)
        
        if not user_order and not delete_pages:
            return Response({'error': 'Either user_order or delete_pages must be provided.'}, status=status.HTTP_400_BAD_REQUEST)
            
        try:
            reader = PdfReader(input_pdf)
            total_pages = len(reader.pages)
            writer = PdfWriter()
            
            if delete_pages:
                # Parse delete pages
                try:
                    if isinstance(delete_pages, str):
                        pages_to_delete = set(map(int, delete_pages.strip('[]').split(',')))
                    else:
                        pages_to_delete = set(map(int, delete_pages))
                except (ValueError, TypeError):
                    return Response({'error': 'Invalid delete pages format.'}, status=status.HTTP_400_BAD_REQUEST)
                
                # Add all pages except deleted ones
                for page_num in range(1, total_pages + 1):
                    if page_num not in pages_to_delete:
                        writer.add_page(reader.pages[page_num - 1])
                        
            elif user_order:
                # Parse user order
                try:
                    if isinstance(user_order, str):
                        user_order = list(map(int, user_order.strip('[]').split(',')))
                    elif isinstance(user_order, list):
                        user_order = list(map(int, user_order))
                except (ValueError, TypeError):
                    return Response({'error': 'Invalid user order format.'}, status=status.HTTP_400_BAD_REQUEST)
                
                # Check if the user's order is valid
                if sorted(user_order) != list(range(1, total_pages + 1)):
                    return Response({'error': 'Invalid page order. Please enter a valid order.'}, status=status.HTTP_400_BAD_REQUEST)
                
                # Add pages in specified order
                for page_number in user_order:
                    writer.add_page(reader.pages[page_number - 1])
            
            # Save organized PDF
            from io import BytesIO
            buffer = BytesIO()
            writer.write(buffer)
            buffer.seek(0)
            
            organized_pdf = OrganizedPdf(user=request.user)
            organized_pdf.organize_pdf.save('organized_output.pdf', ContentFile(buffer.getvalue()))
            organized_pdf.save()
            
            # Get full URL
            protocol = 'https' if request.is_secure() else 'http'
            file_url = f'{protocol}://{request.get_host()}{organized_pdf.organize_pdf.url}'
            
            action = 'deleted' if delete_pages else 'organized'
            response_data = {
                'message': f'PDF pages {action} successfully.',
                'organized_data': {
                    'id': organized_pdf.id,
                    'user': request.user.id,
                    'created_at': organized_pdf.created_at.isoformat(),
                    'organize_pdf': file_url
                }
            }
            return Response(response_data)
            
        except Exception as e:
            return Response({'error': f'An error occurred: {str(e)}'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)



@extend_schema(tags=['PDF Operations'])
class DownloadOrganizedPDFView(APIView):
    permission_classes = [AllowAny]

    def get(self, request, pdf_id, format=None):
        try:
            organized_pdf = OrganizedPdf.objects.get(id=pdf_id)
            return FileResponse(
                organized_pdf.organize_pdf.open('rb'),
                as_attachment=True,
                filename=f'organized_{organized_pdf.id}.pdf'
            )
        except OrganizedPdf.DoesNotExist:
            return Response({'error': 'Organized PDF not found'}, status=404)


@extend_schema(tags=['PDF Operations'])
class UnlockPDFView(APIView):
    permission_classes = [AllowAny]
    parser_classes = (MultiPartParser, FormParser)
    serializer_class = UnlockPDFRequestSerializer

    @extend_schema(
        request=UnlockPDFRequestSerializer,
        responses=inline_serializer(
            name='UnlockPDFResponse',
            fields={
                'message': drf_serializers.CharField(),
                'unlocked_pdf': inline_serializer(
                    name='UnlockPDFData',
                    fields={
                        'id': drf_serializers.IntegerField(),
                        'user': drf_serializers.IntegerField(),
                        'created_at': drf_serializers.DateTimeField(),
                        'unlock_pdf': drf_serializers.URLField(),
                    }
                ),
            }
        ),
        examples=[
            OpenApiExample(
                'Unlock PDF Example',
                value={
                    'message': 'PDF unlocked successfully.',
                    'unlocked_pdf': {
                        'id': 1,
                        'user': 5,
                        'created_at': '2024-01-01T00:00:00Z',
                        'unlock_pdf': 'http://localhost:8000/media/unlocked/output.pdf'
                    }
                }
            )
        ]
    )
    def post(self, request, format=None):
        input_pdf = request.FILES.get('input_pdf', None)
        password = request.data.get('password', '')

        if not input_pdf:
            return Response({'error': 'No input PDF file provided.'}, status=status.HTTP_400_BAD_REQUEST)
        
        if not password:
            return Response({'error': 'No password provided.'}, status=status.HTTP_400_BAD_REQUEST)

        try:
            # Create user if not authenticated
            if request.user.is_authenticated:
                user = request.user
            else:
                from accounts.models import User
                user, created = User.objects.get_or_create(
                    email='test@example.com',
                    defaults={'password': make_password('testpass123')}
                )
            
            # Unlock PDF
            reader = PdfReader(input_pdf)
            if reader.decrypt(password):
                writer = PdfWriter()
                for page in reader.pages:
                    writer.add_page(page)
                
                # Save unlocked PDF
                from io import BytesIO
                buffer = BytesIO()
                writer.write(buffer)
                buffer.seek(0)
                
                unlocked_pdf = UnlockPdf(user=user)
                unlocked_pdf.unlock_pdf.save(
                    f'unlocked_{input_pdf.name}',
                    ContentFile(buffer.getvalue())
                )
                unlocked_pdf.save()
                
                # Get full URL
                protocol = 'https' if request.is_secure() else 'http'
                file_url = f'{protocol}://{request.get_host()}{unlocked_pdf.unlock_pdf.url}'
                
                # Use download endpoint instead of direct media URL
                download_url = f'{protocol}://{request.get_host()}/pdf/download_unlocked_pdf/{unlocked_pdf.id}/'
                
                response_data = {
                    'message': 'PDF unlocked successfully.',
                    'unlocked_pdf': {
                        'id': unlocked_pdf.id,
                        'user': user.id,
                        'created_at': unlocked_pdf.created_at.isoformat(),
                        'unlock_pdf': download_url
                    }
                }
                return Response(response_data)
            else:
                return Response({'error': 'Invalid password'}, status=400)
        except Exception as e:
            return Response({'error': f'An error occurred: {str(e)}'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)




@extend_schema(tags=['PDF Operations'])
class DownloadUnlockedPDFView(APIView):
    permission_classes = [AllowAny]

    @extend_schema(
        responses={
            200: inline_serializer(
                name='UnlockedPDFDownloadResponse',
                fields={
                    'file': drf_serializers.FileField()
                }
            ),
            404: inline_serializer(
                name='UnlockedPDFNotFoundResponse',
                fields={
                    'error': drf_serializers.CharField()
                }
            )
        },
        examples=[
            OpenApiExample(
                'Download Unlocked PDF',
                description='Download the unlocked PDF file',
                response_only=True
            )
        ]
    )
    def get(self, request, pdf_id, format=None):
        try:
            unlocked_pdf = UnlockPdf.objects.get(id=pdf_id)
            
            # Check if file exists
            if not unlocked_pdf.unlock_pdf or not unlocked_pdf.unlock_pdf.name:
                return Response({'error': 'PDF file not found'}, status=404)
            
            try:
                return FileResponse(
                    unlocked_pdf.unlock_pdf.open('rb'),
                    as_attachment=True,
                    filename=f'unlocked_{unlocked_pdf.id}.pdf'
                )
            except FileNotFoundError:
                return Response({'error': 'PDF file not found on disk'}, status=404)
            except Exception as e:
                return Response({'error': f'Error accessing file: {str(e)}'}, status=500)
                
        except UnlockPdf.DoesNotExist:
            return Response({'error': 'Unlocked PDF not found'}, status=404)


@extend_schema(exclude=True)
class UnlockPDFDeleteView(generics.DestroyAPIView):
    queryset = UnlockPdf.objects.all()
    serializer_class = UnlockPdfSerializer
    permission_classes = [IsAuthenticated]



@extend_schema(tags=['PDF Operations'])
class StampPDFView(APIView):
    permission_classes = [IsAuthenticated]
    parser_classes = (MultiPartParser, FormParser)
    serializer_class = StampPDFRequestSerializer

    @extend_schema(
        request=StampPDFRequestSerializer,
        responses=inline_serializer(
            name='StampPDFResponse',
            fields={
                'message': drf_serializers.CharField(),
                'data': inline_serializer(
                    name='StampPDFData',
                    fields={
                        'id': drf_serializers.IntegerField(),
                        'user': drf_serializers.IntegerField(),
                        'created_at': drf_serializers.DateTimeField(),
                        'pdf': drf_serializers.URLField(),
                    }
                ),
            }
        ),
        examples=[
            OpenApiExample(
                'Stamp PDF Example',
                value={
                    'message': 'PDF pages stamped successfully.',
                    'data': {
                        'id': 1,
                        'user': 5,
                        'created_at': '2024-01-01T00:00:00Z',
                        'pdf': 'http://localhost:8000/media/stamped/file.pdf'
                    }
                }
            )
        ]
    )
    def post(self, request, format=None):
        input_pdf = request.FILES.get('input_pdf', None)
        text = request.data.get('text', '')

        if not input_pdf:
            return Response({'error': 'No input PDF file.'}, status=status.HTTP_400_BAD_REQUEST)

        if not text:
            return Response({'error': 'No stamp text provided.'}, status=status.HTTP_400_BAD_REQUEST)

        try:
            new_file = stamp_pdf_with_text(input_pdf, text, request.user)

            serializer = StampPdfSerializer(new_file, context={'request': request})
            return Response({'message': 'PDF pages stamped successfully.', 'data': serializer.data})
        except Exception as e:
            return Response({'error': f'An error occurred: {str(e)}'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@extend_schema(tags=['PDF Operations'])
class SignPDFView(APIView):
    permission_classes = [IsAuthenticated]
    parser_classes = (MultiPartParser, FormParser)

    def post(self, request, format=None):
        input_pdf = request.FILES.get('input_pdf', None)
        signatures = request.data.get('signatures', '[]')

        if not input_pdf:
            return Response({'error': 'No input PDF file provided.'}, status=status.HTTP_400_BAD_REQUEST)

        try:
            # Create user if not authenticated
            if request.user.is_authenticated:
                user = request.user
            else:
                from accounts.models import User
                user, created = User.objects.get_or_create(
                    email='test@example.com',
                    defaults={'password': make_password('testpass123')}
                )
            
            # Sign PDF (simple copy for now)
            reader = PdfReader(input_pdf)
            writer = PdfWriter()
            
            for page in reader.pages:
                writer.add_page(page)
            
            # Save signed PDF
            from io import BytesIO
            buffer = BytesIO()
            writer.write(buffer)
            buffer.seek(0)
            
            # Create a simple model for signed PDF (reuse existing model)
            from .models import StampPdf
            signed_pdf = StampPdf(user=user)
            signed_pdf.pdf.save(
                f'signed_{input_pdf.name}',
                ContentFile(buffer.getvalue())
            )
            signed_pdf.save()
            
            # Get full URL
            protocol = 'https' if request.is_secure() else 'http'
            file_url = f'{protocol}://{request.get_host()}{signed_pdf.pdf.url}'
            
            response_data = {
                'message': 'PDF signing completed.',
                'signed_pdf': {
                    'id': signed_pdf.id,
                    'user': user.id,
                    'created_at': signed_pdf.created_at.isoformat(),
                    'signed_file': file_url
                },
            }
            return Response(response_data)

        except Exception as e:
            return Response({'error': f'An error occurred: {str(e)}'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@extend_schema(tags=['PDF Operations'])
class DownloadOcrPDFView(APIView):
    permission_classes = [AllowAny]

    def get(self, request, pdf_id, format=None):
        try:
            ocr_pdf = OcrPdf.objects.get(id=pdf_id)
            return FileResponse(
                ocr_pdf.pdf.open('rb'),
                as_attachment=True,
                filename=f'ocr_{ocr_pdf.id}.pdf'
            )
        except OcrPdf.DoesNotExist:
            return Response({'error': 'OCR PDF not found'}, status=404)


@extend_schema(tags=['PDF Operations'])
class OcrPDFView(APIView):
    permission_classes = [IsAuthenticated]
    parser_classes = (MultiPartParser, FormParser)
    serializer_class = OcrPDFRequestSerializer

    @extend_schema(
        request=OcrPDFRequestSerializer,
        responses=inline_serializer(
            name='OcrPDFResponse',
            fields={
                'message': drf_serializers.CharField(),
                'data': inline_serializer(
                    name='OcrPDFData',
                    fields={
                        'id': drf_serializers.IntegerField(),
                        'user': drf_serializers.IntegerField(),
                        'language': drf_serializers.CharField(),
                        'created_at': drf_serializers.DateTimeField(),
                        'pdf': drf_serializers.URLField(),
                    }
                ),
                'extracted_text_preview': drf_serializers.CharField(),
                'total_pages_processed': drf_serializers.IntegerField(),
            }
        ),
        examples=[
            OpenApiExample(
                'OCR PDF Example - English',
                value={
                    'message': 'OCR processing completed successfully.',
                    'data': {
                        'id': 1,
                        'user': 5,
                        'language': 'eng',
                        'created_at': '2024-01-01T00:00:00Z',
                        'pdf': 'http://localhost:8000/media/ocr/output.pdf'
                    },
                    'extracted_text_preview': 'Sample extracted text from the PDF...',
                    'total_pages_processed': 5
                }
            ),
            OpenApiExample(
                'OCR PDF Example - Spanish',
                value={
                    'message': 'OCR processing completed successfully.',
                    'data': {
                        'id': 2,
                        'user': 5,
                        'language': 'spa',
                        'created_at': '2024-01-01T00:00:00Z',
                        'pdf': 'http://localhost:8000/media/ocr/output.pdf'
                    },
                    'extracted_text_preview': 'Texto extraído de muestra del PDF...',
                    'total_pages_processed': 3
                }
            )
        ]
    )
    def post(self, request, format=None):
        input_pdf = request.FILES.get('input_pdf', None)
        language = request.data.get('language', 'eng')

        if not input_pdf:
            return Response({'error': 'No input PDF file.'}, status=status.HTTP_400_BAD_REQUEST)

        try:
            user = request.user
            
            # Process PDF with OCR and save to database
            ocr_pdf, _ = pdf_to_ocr(input_pdf, user, language)
            
            serializer = OcrPdfSerializer(ocr_pdf, context={'request': request})
            
            response_data = {
                'message': 'OCR processing completed successfully.',
                'data': serializer.data
            }
            
            return Response(response_data)
            
        except Exception as e:
            return Response({'error': f'OCR processing failed: {str(e)}'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@extend_schema(tags=['PDF Operations'])
class ExtractTextView(APIView):
    permission_classes = [IsAuthenticated]
    parser_classes = (MultiPartParser, FormParser)

    @extend_schema(
        request=inline_serializer(
            name='ExtractTextRequest',
            fields={
                'input_pdf': drf_serializers.FileField(),
                'language': drf_serializers.ChoiceField(
                    choices=[('eng', 'English'), ('spa', 'Spanish'), ('fra', 'French'), ('deu', 'German')],
                    default='eng'
                )
            }
        ),
        responses=inline_serializer(
            name='ExtractTextResponse',
            fields={
                'message': drf_serializers.CharField(),
                'extracted_text': drf_serializers.CharField(),
                'total_pages': drf_serializers.IntegerField(),
                'has_existing_text': drf_serializers.BooleanField(),
            }
        )
    )
    def post(self, request, format=None):
        input_pdf = request.FILES.get('input_pdf', None)
        language = request.data.get('language', 'eng')

        if not input_pdf:
            return Response({'error': 'No input PDF file provided.'}, status=status.HTTP_400_BAD_REQUEST)

        try:
            import tempfile
            from .ocr_processor import PDFOCRProcessor
            
            # Save uploaded file temporarily
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
                for chunk in input_pdf.chunks():
                    temp_file.write(chunk)
                temp_path = temp_file.name
            
            # Extract text using OCR processor
            processor = PDFOCRProcessor()
            full_text = processor.extract_text_only(temp_path, language)
            
            # Get detailed results for metadata
            ocr_results = processor.process_pdf(temp_path, output_format='structured', language=language)
            
            # Clean up
            os.remove(temp_path)
            
            return Response({
                'message': 'Text extraction completed successfully.',
                'extracted_text': full_text,
                'total_pages': ocr_results['total_pages'],
                'has_existing_text': ocr_results['pages_with_existing_text'] > 0,
                'pages_with_existing_text': ocr_results['pages_with_existing_text'],
                'pages_requiring_ocr': ocr_results['pages_requiring_ocr']
            })
            
        except Exception as e:
            return Response({'error': f'Text extraction failed: {str(e)}'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@extend_schema(tags=['PDF Operations'])
class PDFToFormatView(APIView):
    permission_classes = [IsAuthenticated]
    parser_classes = (MultiPartParser, FormParser)
    serializer_class = PDFToFormatRequestSerializer

    @extend_schema(
        request=PDFToFormatRequestSerializer,
        responses=inline_serializer(
            name='PDFToFormatResponse',
            fields={
                'message': drf_serializers.CharField(),
                'conversion_data': inline_serializer(
                    name='PDFToFormatData',
                    fields={
                        'id': drf_serializers.IntegerField(),
                        'user': drf_serializers.IntegerField(),
                        'output_format': drf_serializers.CharField(),
                        'converted_file': drf_serializers.URLField(),
                        'created_at': drf_serializers.DateTimeField(),
                    }
                ),
            }
        )
    )
    def post(self, request, format=None):
        input_pdf = request.FILES.get('input_pdf')
        output_format = request.data.get('output_format')

        if not input_pdf:
            return Response({'error': 'No input PDF file provided.'}, status=status.HTTP_400_BAD_REQUEST)
        
        if not output_format:
            return Response({'error': 'Output format is required.'}, status=status.HTTP_400_BAD_REQUEST)

        try:
            user = request.user
            
            # Convert based on format
            if output_format == 'text':
                converted_content = self._convert_to_text(input_pdf)
                file_extension = 'txt'
            elif output_format in ['jpeg', 'png']:
                converted_content = self._convert_to_images(input_pdf, output_format)
                file_extension = 'zip'
            elif output_format == 'word':
                converted_content = self._convert_to_word(input_pdf)
                file_extension = 'docx'
            elif output_format == 'excel':
                converted_content = self._convert_to_excel(input_pdf)
                file_extension = 'xlsx'
            elif output_format == 'powerpoint':
                converted_content = self._convert_to_powerpoint(input_pdf)
                file_extension = 'pptx'
            else:
                return Response({'error': 'Unsupported output format.'}, status=status.HTTP_400_BAD_REQUEST)
            
            # Save conversion
            conversion = PDFFormatConversion(user=user, output_format=output_format)
            conversion.converted_file.save(
                f'converted_{output_format}.{file_extension}',
                ContentFile(converted_content)
            )
            conversion.save()
            
            # Use download endpoint instead of direct media URL
            protocol = 'https' if request.is_secure() else 'http'
            download_url = f'{protocol}://{request.get_host()}/pdf/download_format_converted/{conversion.id}/'
            
            response_data = {
                'message': f'PDF to {output_format} conversion completed.',
                'conversion_data': {
                    'id': conversion.id,
                    'user': user.id,
                    'output_format': output_format,
                    'converted_file': download_url,
                    'created_at': conversion.created_at.isoformat()
                }
            }
            return Response(response_data)
            
        except Exception as e:
            return Response({'error': f'Conversion failed: {str(e)}'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
    
    def _convert_to_text(self, pdf_file):
        try:
            import tempfile
            try:
                import pdfplumber
                return self._convert_to_text_advanced(pdf_file)
            except ImportError:
                return self._convert_to_text_basic(pdf_file)
        except Exception as e:
            return f'Error extracting text from PDF: {str(e)}'.encode('utf-8')
    
    def _convert_to_text_advanced(self, pdf_file):
        import tempfile
        import pdfplumber
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
            pdf_file.seek(0)
            for chunk in pdf_file.chunks():
                temp_file.write(chunk)
            temp_path = temp_file.name
        
        text_parts = []
        with pdfplumber.open(temp_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text(layout=True)
                if page_text and page_text.strip():
                    text_parts.append(f"--- Page {page_num} ---\n{page_text.strip()}\n")
                
                tables = page.extract_tables()
                for i, table in enumerate(tables):
                    text_parts.append(f"\n--- Table {i+1} on Page {page_num} ---\n")
                    for row in table:
                        if row:
                            text_parts.append('\t'.join([str(cell) if cell else '' for cell in row]))
                    text_parts.append('\n')
        
        os.remove(temp_path)
        final_text = '\n'.join(text_parts) if text_parts else 'No text content found in PDF'
        return final_text.encode('utf-8')
    
    def _convert_to_text_basic(self, pdf_file):
        from PyPDF2 import PdfReader
        pdf_file.seek(0)
        reader = PdfReader(pdf_file)
        
        text_parts = []
        for page_num, page in enumerate(reader.pages, 1):
            page_text = page.extract_text()
            if page_text.strip():
                text_parts.append(f"--- Page {page_num} ---\n{page_text.strip()}\n")
        
        final_text = '\n'.join(text_parts) if text_parts else 'No text content found in PDF'
        return final_text.encode('utf-8')
    
    def _convert_to_images(self, pdf_file, format_type):
        # Reuse existing PDF to image logic
        from .utils import convert_pdf_to_image, create_zip_file
        images = convert_pdf_to_image(pdf_file, format_type)
        _, zip_content = create_zip_file(images, None, format_type)
        return zip_content
    
    def _convert_to_word(self, pdf_file):
        try:
            import fitz
            from docx import Document
            from docx.shared import Inches, Pt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.enum.style import WD_STYLE_TYPE
            import tempfile
            from io import BytesIO
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
                pdf_file.seek(0)
                for chunk in pdf_file.chunks():
                    temp_file.write(chunk)
                temp_path = temp_file.name
            
            doc = Document()
            pdf_doc = fitz.open(temp_path)
            
            for page_num in range(len(pdf_doc)):
                if page_num > 0:
                    doc.add_page_break()
                
                page = pdf_doc[page_num]
                
                # Extract text blocks with formatting
                blocks = page.get_text("dict")
                
                for block in blocks["blocks"]:
                    if "lines" in block:
                        for line in block["lines"]:
                            paragraph = doc.add_paragraph()
                            
                            for span in line["spans"]:
                                text = span["text"]
                                if not text.strip():
                                    continue
                                
                                run = paragraph.add_run(text)
                                
                                # Extract font information
                                font_name = span.get("font", "Times New Roman")
                                font_size = span.get("size", 12)
                                font_flags = span.get("flags", 0)
                                font_color = span.get("color", 0)
                                
                                # Apply font name
                                if font_name:
                                    run.font.name = font_name
                                
                                # Apply font size (reduced by 20%)
                                run.font.size = Pt(max(8, font_size * 0.8))
                                
                                # Apply font styles (correct PyMuPDF flags)
                                if font_flags & 16:  # Bold (2^4)
                                    run.font.bold = True
                                if font_flags & 2:   # Italic (2^1)
                                    run.font.italic = True
                                # Only apply underline if explicitly detected
                                # PyMuPDF doesn't reliably detect underline via flags
                                # Skip underline to avoid false positives
                                
                                # Apply font color
                                if font_color != 0:
                                    r = (font_color >> 16) & 255
                                    g = (font_color >> 8) & 255
                                    b = font_color & 255
                                    run.font.color.rgb = RGBColor(r, g, b)
                            
                            # Set paragraph alignment based on text position
                            bbox = line["bbox"]
                            page_width = page.rect.width
                            left_margin = bbox[0]
                            right_margin = page_width - bbox[2]
                            text_width = bbox[2] - bbox[0]
                            
                            # Detect justified text (spans most of page width)
                            if text_width > page_width * 0.7 and left_margin < page_width * 0.15 and right_margin < page_width * 0.15:
                                paragraph.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                            elif left_margin < page_width * 0.1:  # Text starts very close to left
                                paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
                            elif right_margin < page_width * 0.1:  # Text ends very close to right
                                paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                            elif abs(left_margin - right_margin) < page_width * 0.03:  # Nearly equal margins
                                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            else:
                                paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT  # Default to left
                
                # Extract and add images
                image_list = page.get_images()
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        base_image = pdf_doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as img_temp:
                            img_temp.write(image_bytes)
                            img_temp_path = img_temp.name
                        
                        try:
                            doc.add_picture(img_temp_path, width=Inches(6))
                        except:
                            pass
                        finally:
                            os.remove(img_temp_path)
                    except:
                        continue
            
            pdf_doc.close()
            os.remove(temp_path)
            
            buffer = BytesIO()
            doc.save(buffer)
            return buffer.getvalue()
            
        except Exception as e:
            return self._convert_to_word_fallback(pdf_file)
    
    def _convert_to_word_fallback(self, pdf_file):
        from docx import Document
        from docx.shared import Inches
        from io import BytesIO
        import tempfile
        try:
            import fitz
            import pdfplumber
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
                pdf_file.seek(0)
                for chunk in pdf_file.chunks():
                    temp_file.write(chunk)
                temp_path = temp_file.name
            
            doc = Document()
            pdf_doc = fitz.open(temp_path)
            
            with pdfplumber.open(temp_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    if page_num > 0:
                        doc.add_page_break()
                    
                    # Extract images
                    try:
                        fitz_page = pdf_doc[page_num]
                        image_list = fitz_page.get_images()
                        
                        for img in image_list:
                            try:
                                xref = img[0]
                                base_image = pdf_doc.extract_image(xref)
                                image_bytes = base_image["image"]
                                
                                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as img_temp:
                                    img_temp.write(image_bytes)
                                    img_temp_path = img_temp.name
                                
                                try:
                                    doc.add_picture(img_temp_path, width=Inches(6))
                                except Exception:
                                    pass
                                finally:
                                    os.remove(img_temp_path)
                            except Exception:
                                continue
                    except Exception:
                        pass
                    
                    # Extract text with layout preservation
                    page_text = page.extract_text(layout=True)
                    if page_text and page_text.strip():
                        lines = page_text.split('\n')
                        for line in lines:
                            if line.strip():
                                doc.add_paragraph(line.strip())
            
            pdf_doc.close()
            os.remove(temp_path)
            
            buffer = BytesIO()
            doc.save(buffer)
            return buffer.getvalue()
        except ImportError:
            return self._convert_to_word_basic(pdf_file)
        except Exception:
            return self._convert_to_word_basic(pdf_file)
    
    def _convert_to_word_professional(self, pdf_file):
        from pdf2docx import Converter
        from io import BytesIO
        import tempfile
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            pdf_file.seek(0)
            for chunk in pdf_file.chunks():
                temp_pdf.write(chunk)
            temp_pdf_path = temp_pdf.name
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_docx:
            temp_docx_path = temp_docx.name
        
        cv = Converter(temp_pdf_path)
        cv.convert(temp_docx_path, start=0, end=None)
        cv.close()
        
        with open(temp_docx_path, 'rb') as f:
            docx_content = f.read()
        
        os.remove(temp_pdf_path)
        os.remove(temp_docx_path)
        
        return docx_content
    
    def _convert_to_word_advanced(self, pdf_file):
        from docx import Document
        from docx.shared import Inches
        from io import BytesIO
        import tempfile
        import pdfplumber
        import fitz
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
            pdf_file.seek(0)
            for chunk in pdf_file.chunks():
                temp_file.write(chunk)
            temp_path = temp_file.name
        
        doc = Document()
        
        try:
            pdf_doc = fitz.open(temp_path)
            
            with pdfplumber.open(temp_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    if page_num > 0:
                        doc.add_page_break()
                    
                    # Extract images
                    try:
                        fitz_page = pdf_doc[page_num]
                        image_list = fitz_page.get_images()
                        
                        for img in image_list:
                            try:
                                xref = img[0]
                                base_image = pdf_doc.extract_image(xref)
                                image_bytes = base_image["image"]
                                
                                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as img_temp:
                                    img_temp.write(image_bytes)
                                    img_temp_path = img_temp.name
                                
                                try:
                                    doc.add_picture(img_temp_path, width=Inches(6))
                                except Exception:
                                    pass
                                finally:
                                    os.remove(img_temp_path)
                            except Exception:
                                continue
                    except Exception:
                        pass
                    
                    # Extract text with layout preservation
                    page_text = page.extract_text(layout=True)
                    if page_text and page_text.strip():
                        lines = page_text.split('\n')
                        for line in lines:
                            if line.strip():
                                doc.add_paragraph(line.strip())
                    
                    # Add tables
                    tables = page.extract_tables()
                    for table_data in tables:
                        if table_data and len(table_data) > 0 and len(table_data[0]) > 0:
                            table = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
                            table.style = 'Table Grid'
                            for i, row in enumerate(table_data):
                                for j, cell in enumerate(row):
                                    if cell and j < len(table.columns):
                                        table.cell(i, j).text = str(cell)
            
            pdf_doc.close()
            
        except Exception:
            # Fallback
            with pdfplumber.open(temp_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    if page_num > 0:
                        doc.add_page_break()
                    
                    page_text = page.extract_text(layout=True)
                    if page_text and page_text.strip():
                        lines = page_text.split('\n')
                        for line in lines:
                            if line.strip():
                                doc.add_paragraph(line.strip())
        
        os.remove(temp_path)
        buffer = BytesIO()
        doc.save(buffer)
        return buffer.getvalue()
    
    def _convert_to_word_basic(self, pdf_file):
        from docx import Document
        from docx.shared import Inches
        from io import BytesIO
        from PyPDF2 import PdfReader
        import tempfile
        import fitz
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
            pdf_file.seek(0)
            for chunk in pdf_file.chunks():
                temp_file.write(chunk)
            temp_path = temp_file.name
        
        doc = Document()
        pdf_file.seek(0)
        reader = PdfReader(pdf_file)
        
        try:
            pdf_doc = fitz.open(temp_path)
            
            for page_num, page in enumerate(reader.pages):
                if page_num > 0:
                    doc.add_page_break()
                
                # Extract images
                try:
                    fitz_page = pdf_doc[page_num]
                    image_list = fitz_page.get_images()
                    
                    for img in image_list:
                        try:
                            xref = img[0]
                            base_image = pdf_doc.extract_image(xref)
                            image_bytes = base_image["image"]
                            
                            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as img_temp:
                                img_temp.write(image_bytes)
                                img_temp_path = img_temp.name
                            
                            try:
                                doc.add_picture(img_temp_path, width=Inches(6))
                            except Exception:
                                pass
                            finally:
                                os.remove(img_temp_path)
                        except Exception:
                            continue
                except Exception:
                    pass
                
                # Extract text preserving original formatting
                page_text = page.extract_text()
                if page_text.strip():
                    lines = page_text.split('\n')
                    for line in lines:
                        if line.strip():
                            doc.add_paragraph(line.strip())
            
            pdf_doc.close()
            
        except Exception:
            # Fallback without images
            for page_num, page in enumerate(reader.pages):
                if page_num > 0:
                    doc.add_page_break()
                
                page_text = page.extract_text()
                if page_text.strip():
                    lines = page_text.split('\n')
                    for line in lines:
                        if line.strip():
                            doc.add_paragraph(line.strip())
        
        os.remove(temp_path)
        buffer = BytesIO()
        doc.save(buffer)
        return buffer.getvalue()
    
    def _convert_to_excel(self, pdf_file):
        # Basic text to Excel
        import pandas as pd
        from io import BytesIO
        
        text = self._convert_to_text(pdf_file).decode('utf-8')
        df = pd.DataFrame({'Content': [text]})
        
        buffer = BytesIO()
        df.to_excel(buffer, index=False)
        return buffer.getvalue()
    
    def _convert_to_powerpoint(self, pdf_file):
        try:
            import pdfplumber
            import fitz
            return self._convert_to_powerpoint_advanced(pdf_file)
        except ImportError:
            return self._convert_to_powerpoint_basic(pdf_file)
        except Exception:
            return self._convert_to_powerpoint_basic(pdf_file)
    
    def _convert_to_powerpoint_advanced(self, pdf_file):
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from io import BytesIO
        import tempfile
        import pdfplumber
        import fitz
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
            pdf_file.seek(0)
            for chunk in pdf_file.chunks():
                temp_file.write(chunk)
            temp_path = temp_file.name
        
        prs = Presentation()
        pdf_doc = fitz.open(temp_path)
        
        for page_num in range(len(pdf_doc)):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            page = pdf_doc[page_num]
            
            # Convert entire page to image first
            mat = fitz.Matrix(2.0, 2.0)  # High resolution
            pix = page.get_pixmap(matrix=mat)
            img_data = pix.tobytes("png")
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as page_img:
                page_img.write(img_data)
                page_img_path = page_img.name
            
            try:
                # Add page as image
                slide.shapes.add_picture(
                    page_img_path,
                    Inches(0.5),
                    Inches(0.5),
                    width=Inches(9),
                    height=Inches(6.5)
                )
            except Exception as e:
                print(f"Error adding page image: {e}")
            finally:
                try:
                    os.remove(page_img_path)
                except:
                    pass
            
            # Add text below image with very small font
            try:
                with pdfplumber.open(temp_path) as pdf_plumber:
                    plumber_page = pdf_plumber.pages[page_num]
                    page_text = plumber_page.extract_text()
                    
                    if page_text and page_text.strip():
                        textbox = slide.shapes.add_textbox(
                            Inches(0.5),
                            Inches(7.2),
                            Inches(9),
                            Inches(0.8)
                        )
                        text_frame = textbox.text_frame
                        text_frame.word_wrap = True
                        
                        clean_text = page_text.strip()[:300]
                        if len(page_text) > 300:
                            clean_text += '...'
                        
                        text_frame.text = clean_text
                        
                        for paragraph in text_frame.paragraphs:
                            paragraph.font.size = Pt(4)
                            paragraph.font.name = 'Arial'
            except Exception as e:
                print(f"Error adding text: {e}")
        
        pdf_doc.close()
        os.remove(temp_path)
        
        buffer = BytesIO()
        prs.save(buffer)
        return buffer.getvalue()
    
    def _convert_to_powerpoint_basic(self, pdf_file):
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from io import BytesIO
        import tempfile
        import fitz
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
            pdf_file.seek(0)
            for chunk in pdf_file.chunks():
                temp_file.write(chunk)
            temp_path = temp_file.name
        
        prs = Presentation()
        pdf_doc = fitz.open(temp_path)
        
        for page_num in range(len(pdf_doc)):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            page = pdf_doc[page_num]
            
            # Convert page to image
            pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5))
            img_data = pix.tobytes("png")
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as img_temp:
                img_temp.write(img_data)
                img_temp_path = img_temp.name
            
            try:
                slide.shapes.add_picture(img_temp_path, Inches(1), Inches(1), Inches(8), Inches(6))
            except:
                pass
            finally:
                try:
                    os.remove(img_temp_path)
                except:
                    pass
        
        pdf_doc.close()
        os.remove(temp_path)
        
        buffer = BytesIO()
        prs.save(buffer)
        return buffer.getvalue()


@extend_schema(tags=['PDF Operations'])
class FormatToPDFView(APIView):
    permission_classes = [IsAuthenticated]
    parser_classes = (MultiPartParser, FormParser)
    serializer_class = WordToPdfRequestSerializer

    @extend_schema(
        request=WordToPdfRequestSerializer,
        responses=inline_serializer(
            name='FormatToPDFResponse',
            fields={
                'message': drf_serializers.CharField(),
                'conversion_data': inline_serializer(
                    name='FormatToPDFData',
                    fields={
                        'id': drf_serializers.IntegerField(),
                        'user': drf_serializers.IntegerField(),
                        'input_format': drf_serializers.CharField(),
                        'converted_pdf': drf_serializers.URLField(),
                        'created_at': drf_serializers.DateTimeField(),
                    }
                ),
            }
        )
    )
    def post(self, request, format=None):
        input_file = request.FILES.get('input_file')
        input_format = request.data.get('input_format')

        if not input_file:
            return Response({'error': 'No input file provided.'}, status=status.HTTP_400_BAD_REQUEST)
        
        if not input_format:
            file_ext = input_file.name.split('.')[-1].lower()
            input_format = file_ext

        try:
            user = request.user
            
            if input_format in ['doc', 'docx']:
                pdf_content = self._convert_word_to_pdf(input_file)
            elif input_format in ['xls', 'xlsx']:
                pdf_content = self._convert_excel_to_pdf(input_file)
            elif input_format == 'csv':
                pdf_content = self._convert_csv_to_pdf(input_file)
            elif input_format in ['ppt', 'pptx']:
                pdf_content = self._convert_powerpoint_to_pdf(input_file)
            elif input_format == 'txt':
                pdf_content = self._convert_text_to_pdf(input_file)
            elif input_format == 'rtf':
                pdf_content = self._convert_rtf_to_pdf(input_file)
            elif input_format in ['jpg', 'jpeg', 'png', 'gif', 'bmp']:
                pdf_content = self._convert_image_to_pdf(input_file)
            else:
                return Response({'error': f'Unsupported input format: {input_format}'}, status=status.HTTP_400_BAD_REQUEST)
            
            # Use existing WordToPdfConversion model
            conversion_instance = WordToPdfConversion(user=user)
            conversion_instance.save()
            
            word_to_pdf_instance = WordToPdf()
            word_to_pdf_instance.word_to_pdf.save(
                f'{input_format}_to_pdf.pdf',
                ContentFile(pdf_content)
            )
            word_to_pdf_instance.save()
            
            conversion_instance.word_to_pdfs.add(word_to_pdf_instance)
            conversion_instance.save()
            
            # Use download endpoint instead of direct media URL
            protocol = 'https' if request.is_secure() else 'http'
            download_url = f'{protocol}://{request.get_host()}/pdf/download_converted_pdf/{word_to_pdf_instance.id}/'
            
            response_data = {
                'message': f'{input_format.upper()} to PDF conversion completed.',
                'conversion_data': {
                    'id': conversion_instance.id,
                    'user': user.id,
                    'input_format': input_format,
                    'converted_pdf': download_url,
                    'created_at': conversion_instance.created_at.isoformat()
                }
            }
            return Response(response_data)
            
        except Exception as e:
            return Response({'error': f'Conversion failed: {str(e)}'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
    
    def _convert_word_to_pdf(self, file):
        from docx import Document
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        from io import BytesIO
        
        doc = Document(file)
        text_content = '\n'.join([p.text for p in doc.paragraphs])
        
        buffer = BytesIO()
        p = canvas.Canvas(buffer, pagesize=letter)
        width, height = letter
        
        y = height - 50
        for line in text_content.split('\n'):
            if y < 50:
                p.showPage()
                y = height - 50
            p.drawString(50, y, line[:80])
            y -= 20
        
        p.save()
        return buffer.getvalue()
    
    def _convert_excel_to_pdf(self, file):
        import pandas as pd
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        from io import BytesIO
        
        df = pd.read_excel(file)
        
        buffer = BytesIO()
        p = canvas.Canvas(buffer, pagesize=letter)
        width, height = letter
        
        y = height - 50
        p.drawString(50, y, f'Excel Data ({len(df)} rows, {len(df.columns)} columns)')
        y -= 30
        
        headers = ' | '.join(df.columns[:5])
        p.drawString(50, y, headers)
        y -= 20
        
        for idx, row in df.head(20).iterrows():
            if y < 50:
                p.showPage()
                y = height - 50
            row_text = ' | '.join([str(val)[:15] for val in row.values[:5]])
            p.drawString(50, y, row_text)
            y -= 15
        
        p.save()
        return buffer.getvalue()
    
    def _convert_csv_to_pdf(self, file):
        import pandas as pd
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        from io import BytesIO
        
        df = pd.read_csv(file)
        
        buffer = BytesIO()
        p = canvas.Canvas(buffer, pagesize=letter)
        width, height = letter
        
        y = height - 50
        p.drawString(50, y, f'CSV Data ({len(df)} rows, {len(df.columns)} columns)')
        y -= 30
        
        headers = ' | '.join(df.columns[:5])
        p.drawString(50, y, headers)
        y -= 20
        
        for idx, row in df.head(20).iterrows():
            if y < 50:
                p.showPage()
                y = height - 50
            row_text = ' | '.join([str(val)[:15] for val in row.values[:5]])
            p.drawString(50, y, row_text)
            y -= 15
        
        p.save()
        return buffer.getvalue()
    
    def _convert_powerpoint_to_pdf(self, file):
        from pptx import Presentation
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        from io import BytesIO
        
        prs = Presentation(file)
        
        buffer = BytesIO()
        p = canvas.Canvas(buffer, pagesize=letter)
        width, height = letter
        
        for slide_num, slide in enumerate(prs.slides, 1):
            y = height - 50
            p.drawString(50, y, f'Slide {slide_num}')
            y -= 30
            
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    for line in shape.text.split('\n'):
                        if y < 50:
                            p.showPage()
                            y = height - 50
                        p.drawString(70, y, line[:70])
                        y -= 15
            
            if slide_num < len(prs.slides):
                p.showPage()
        
        p.save()
        return buffer.getvalue()
    
    def _convert_text_to_pdf(self, file):
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        from io import BytesIO
        
        text_content = file.read().decode('utf-8')
        
        buffer = BytesIO()
        p = canvas.Canvas(buffer, pagesize=letter)
        width, height = letter
        
        y = height - 50
        for line in text_content.split('\n'):
            if y < 50:
                p.showPage()
                y = height - 50
            p.drawString(50, y, line[:80])
            y -= 20
        
        p.save()
        return buffer.getvalue()
    
    def _convert_rtf_to_pdf(self, file):
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        from io import BytesIO
        
        try:
            content = file.read().decode('utf-8')
        except:
            content = str(file.read())
        
        buffer = BytesIO()
        p = canvas.Canvas(buffer, pagesize=letter)
        width, height = letter
        
        y = height - 50
        for line in content.split('\n'):
            if y < 50:
                p.showPage()
                y = height - 50
            p.drawString(50, y, line[:80])
            y -= 20
        
        p.save()
        return buffer.getvalue()
    
    def _convert_image_to_pdf(self, file):
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.utils import ImageReader
        from io import BytesIO
        from PIL import Image
        
        image = Image.open(file)
        
        buffer = BytesIO()
        p = canvas.Canvas(buffer, pagesize=letter)
        width, height = letter
        
        img_width, img_height = image.size
        aspect_ratio = img_height / img_width
        
        max_width = width - 100
        max_height = height - 100
        
        if img_width > max_width:
            img_width = max_width
            img_height = img_width * aspect_ratio
        
        if img_height > max_height:
            img_height = max_height
            img_width = img_height / aspect_ratio
        
        x = (width - img_width) / 2
        y = (height - img_height) / 2
        
        file.seek(0)
        img_reader = ImageReader(file)
        p.drawImage(img_reader, x, y, width=img_width, height=img_height)
        
        p.save()
        return buffer.getvalue()


@extend_schema(tags=['PDF Operations'])
class DownloadConvertedPDFView(APIView):
    permission_classes = [AllowAny]

    def get(self, request, pdf_id, format=None):
        try:
            word_to_pdf = WordToPdf.objects.get(id=pdf_id)
            return FileResponse(
                word_to_pdf.word_to_pdf.open('rb'),
                as_attachment=True,
                filename=f'converted_{word_to_pdf.id}.pdf'
            )
        except WordToPdf.DoesNotExist:
            return Response({'error': 'Converted PDF not found'}, status=404)

@extend_schema(tags=['PDF Operations'])
class DownloadFormatConvertedView(APIView):
    permission_classes = [AllowAny]

    def get(self, request, conversion_id, format=None):
        try:
            conversion = PDFFormatConversion.objects.get(id=conversion_id)
            return FileResponse(
                conversion.converted_file.open('rb'),
                as_attachment=True,
                filename=f'converted_{conversion.output_format}_{conversion.id}.{conversion.converted_file.name.split(".")[-1]}'
            )
        except PDFFormatConversion.DoesNotExist:
            return Response({'error': 'Converted file not found'}, status=404)
